from pathlib import Path
from typing import Dict, Any
from .base import BaseConverter
from pptx import Presentation
from pptx.util import Inches

class PPTConverter(BaseConverter):
    SUPPORTED_EXTENSIONS = ['.pptx']
    
    async def convert_to_markdown(self, file_path: Path) -> Dict[str, Any]:
        """PPT转Markdown"""
        try:
            prs = Presentation(str(file_path))
            
            markdown_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                markdown_content.append(f"# Slide {i}")
                markdown_content.append("")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        markdown_content.append(shape.text)
                        markdown_content.append("")
                    
                    # 处理表格
                    if shape.has_table:
                        table = shape.table
                        table_content = self._convert_table_to_markdown(table)
                        markdown_content.append(table_content)
                        markdown_content.append("")
                
                markdown_content.append("---")
                markdown_content.append("")
            
            return {
                "success": True,
                "content": "\n".join(markdown_content),
                "slides": len(prs.slides)
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _convert_table_to_markdown(self, table) -> str:
        """将PPT表格转换为Markdown表格"""
        rows = []
        
        # 处理表头
        if len(table.rows) > 0:
            header = " | ".join(cell.text for cell in table.rows[0].cells)
            rows.append(f"| {header} |")
            rows.append("| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |")
        
        # 处理数据行
        for i, row in enumerate(table.rows):
            if i == 0:  # 跳过表头
                continue
            row_content = " | ".join(cell.text for cell in row.cells)
            rows.append(f"| {row_content} |")
        
        return "\n".join(rows)
    
    async def convert_from_markdown(self, md_content: str, output_path: Path) -> Path:
        """Markdown转PPT"""
        prs = Presentation()
        
        lines = md_content.split('\n')
        current_slide = None
        slide_content = []
        
        for line in lines:
            line = line.strip()
            
            # 检测新幻灯片
            if line.startswith('# '):
                # 保存当前幻灯片内容
                if current_slide is not None and slide_content:
                    self._add_content_to_slide(current_slide, '\n'.join(slide_content))
                
                # 创建新幻灯片
                title = line[2:].strip()
                current_slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
                
                # 设置标题
                title_shape = current_slide.shapes.title
                title_shape.text = title
                
                slide_content = []
            
            elif line == '---':
                # 幻灯片分隔符
                if current_slide is not None and slide_content:
                    self._add_content_to_slide(current_slide, '\n'.join(slide_content))
                
                current_slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
                slide_content = []
            
            elif line:
                slide_content.append(line)
        
        # 处理最后一张幻灯片
        if current_slide is not None and slide_content:
            self._add_content_to_slide(current_slide, '\n'.join(slide_content))
        
        prs.save(str(output_path))
        return output_path
    
    def _add_content_to_slide(self, slide, content: str):
        """向幻灯片添加内容"""
        # 添加文本框
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        for i, paragraph in enumerate(content.split('\n')):
            if i == 0:
                tf.text = paragraph
            else:
                p = tf.add_paragraph()
                p.text = paragraph